#!/usr/bin/env python3
"""Generate docs/DayReminders-Overview.pptx — a 9-slide deck for announcing
Day Reminders to the team.

Each screenshot slot is a real Picture object backed by a generated gray
placeholder PNG, so users can right-click any one and choose "Change
Picture" in PowerPoint to swap in the actual screenshot without resizing.

Re-run after capturing screenshots or after content changes:
    py tools/build-overview-pptx.py
"""

from pathlib import Path
from io import BytesIO

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from PIL import Image, ImageDraw, ImageFont

ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "docs" / "DayReminders-Overview.pptx"
SCREENSHOTS_DIR = ROOT / "docs" / "screenshots"

# Map each slide's screenshot index to the real PNG filename. If the file
# exists in docs/screenshots/, it gets embedded; otherwise a labeled
# placeholder is generated. Slides 7 and 8 are USERGUIDE-only.
SCREENSHOT_FILES = {
    1: "01-left-rail.png",
    2: "02-add-form.png",
    3: "03-lines-view.png",
    4: "04-week-view.png",
    5: "05-group-by-client.png",
    6: "06-proactive-card.png",
}

# App accent + supporting colors
ACCENT = RGBColor(0x38, 0xAE, 0xEB)
INK_DARK = RGBColor(0x18, 0x2C, 0x3C)
TEXT = RGBColor(0x24, 0x24, 0x24)
MUTED = RGBColor(0x61, 0x61, 0x61)


def _load_font(name: str, size: int):
    """Try the named Windows font, fall back to the default if missing."""
    try:
        return ImageFont.truetype(name, size)
    except Exception:
        try:
            return ImageFont.truetype("segoeui.ttf", size)
        except Exception:
            return ImageFont.load_default()


def make_placeholder_png(idx: int, caption: str, w: int = 1600, h: int = 900) -> bytes:
    """Generate a clearly-labeled gray placeholder PNG. Replaced by the user
    via right-click → Change Picture in PowerPoint."""
    img = Image.new("RGB", (w, h), (232, 236, 239))
    draw = ImageDraw.Draw(img)
    # Outer dashed-ish border (drawn as repeated short segments)
    border_color = (176, 184, 192)
    seg = 24
    for x in range(8, w - 8, seg * 2):
        draw.line([(x, 8), (min(x + seg, w - 8), 8)], fill=border_color, width=4)
        draw.line([(x, h - 8), (min(x + seg, w - 8), h - 8)], fill=border_color, width=4)
    for y in range(8, h - 8, seg * 2):
        draw.line([(8, y), (8, min(y + seg, h - 8))], fill=border_color, width=4)
        draw.line([(w - 8, y), (w - 8, min(y + seg, h - 8))], fill=border_color, width=4)
    # Subtle diagonal hatch
    for x in range(-h, w, 80):
        draw.line([(x, 0), (x + h, h)], fill=(216, 222, 228), width=1)

    title_font = _load_font("segoeuib.ttf", 90)
    sub_font = _load_font("segoeui.ttf", 40)

    title_text = f"Screenshot {idx:02d}"
    tb = draw.textbbox((0, 0), title_text, font=title_font)
    tw, th = tb[2] - tb[0], tb[3] - tb[1]
    draw.text(((w - tw) // 2, h // 2 - th - 30), title_text,
              fill=(60, 75, 95), font=title_font)

    # Word-wrap the caption
    max_w = int(w * 0.85)
    words = caption.split()
    lines, cur = [], ""
    for word in words:
        test = (cur + " " + word).strip()
        tw2 = draw.textbbox((0, 0), test, font=sub_font)[2]
        if tw2 <= max_w:
            cur = test
        else:
            if cur:
                lines.append(cur)
            cur = word
    if cur:
        lines.append(cur)
    y = h // 2 + 30
    for line in lines:
        lb = draw.textbbox((0, 0), line, font=sub_font)
        lw, lh = lb[2] - lb[0], lb[3] - lb[1]
        draw.text(((w - lw) // 2, y), line, fill=(85, 101, 117), font=sub_font)
        y += int(lh * 1.4)

    buf = BytesIO()
    img.save(buf, format="PNG", optimize=True)
    return buf.getvalue()


def add_title_bar(slide, prs, title: str, color=ACCENT):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.9))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    tf = bar.text_frame
    tf.margin_left = Inches(0.5)
    tf.margin_right = Inches(0.5)
    tf.margin_top = Inches(0.15)
    tf.margin_bottom = Inches(0.15)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    p.font.name = "Segoe UI"


def add_body_bullets(slide, prs, lines, left, top, width, height):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(17)
        p.font.color.rgb = TEXT
        p.font.name = "Segoe UI"
        p.space_after = Pt(6)


def add_placeholder_picture(slide, prs, idx: int, caption: str,
                            left: Emu, top: Emu, width: Emu, height: Emu):
    # Prefer the real screenshot. Critically: fit it within the slot while
    # preserving its native aspect ratio (centered, letterbox-style). Without
    # this, PowerPoint stretches every screenshot to fill the slot exactly,
    # which is what made the v1 deck look distorted.
    real = SCREENSHOTS_DIR / SCREENSHOT_FILES.get(idx, "")
    if real.exists() and real.stat().st_size > 0:
        with Image.open(real) as im:
            iw, ih = im.size
        img_ar = iw / ih
        slot_ar = width / height
        if img_ar > slot_ar:
            # Image wider than slot — fit to width, center vertically
            w = width
            h = int(width / img_ar)
            l = left
            t = top + (height - h) // 2
        else:
            # Image taller than slot — fit to height, center horizontally
            h = height
            w = int(height * img_ar)
            l = left + (width - w) // 2
            t = top
        return slide.shapes.add_picture(str(real), l, t, w, h)
    # Placeholder fills the slot (the placeholder is a deliberate guide, not a
    # real screenshot; stretching it is fine and keeps the slot fully marked).
    png_bytes = make_placeholder_png(idx, caption)
    return slide.shapes.add_picture(BytesIO(png_bytes), left, top, width, height)


def add_footer(slide, prs, text: str):
    tb = slide.shapes.add_textbox(Inches(0.5), prs.slide_height - Inches(0.45),
                                  prs.slide_width - Inches(1.0), Inches(0.35))
    p = tb.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(10)
    p.font.color.rgb = MUTED
    p.font.name = "Segoe UI"
    p.alignment = PP_ALIGN.RIGHT


def build_deck():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # 1. Title
    s = prs.slides.add_slide(blank)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = INK_DARK
    bg.line.fill.background()
    tb = s.shapes.add_textbox(Inches(1), Inches(2.3), Inches(11), Inches(1.8))
    p = tb.text_frame.paragraphs[0]
    p.text = "Day Reminders"
    p.font.size = Pt(80); p.font.bold = True
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF); p.font.name = "Segoe UI"
    tb = s.shapes.add_textbox(Inches(1), Inches(3.9), Inches(11), Inches(0.8))
    p = tb.text_frame.paragraphs[0]
    p.text = "Reminders, in the place you already work."
    p.font.size = Pt(32); p.font.color.rgb = ACCENT; p.font.name = "Segoe UI"
    tb = s.shapes.add_textbox(Inches(1), Inches(6.4), Inches(11), Inches(0.5))
    p = tb.text_frame.paragraphs[0]
    p.text = "v1.4.6  ·  Kation Technologies  ·  June 2026"
    p.font.size = Pt(14); p.font.color.rgb = RGBColor(0xB0, 0xC8, 0xDC); p.font.name = "Segoe UI"

    # 2. What + Where
    s = prs.slides.add_slide(blank)
    add_title_bar(s, prs, "What it is")
    add_body_bullets(s, prs, [
        "Add what you need to remember today.",
        "The bot pings you in chat before each one.",
        "Mark it done from the card. Snooze it. Move on.",
        "",
        "No flipping to a separate app — it lives where you already work.",
        "",
        "Where to find it:",
        "•  Alarm-clock icon in your Teams left rail.",
        "•  Two top tabs: Reminders (the UI) and Chat (where the bot pings you).",
    ], Inches(0.5), Inches(1.3), Inches(5.3), Inches(5.5))
    add_placeholder_picture(s, prs,1, "Teams left rail with the Day Reminders icon",
                            Inches(6.0), Inches(1.3), Inches(6.83), Inches(5.5))
    add_footer(s, prs, "Day Reminders v1.4.6")

    # 3. Adding from tab
    s = prs.slides.add_slide(blank)
    add_title_bar(s, prs, "Adding a reminder")
    add_body_bullets(s, prs, [
        "Top of the Reminders tab — fill in the add form:",
        "•  Title — what to remember. Add #tags inline (e.g. Send report #work).",
        "•  Client — optional, tracks which engagement; autocompletes from past clients.",
        "•  Due date — defaults to today; pick any other day.",
        "•  Time — optional. Leave blank for an 'anytime today' item.",
        "•  + Details — toggle for notes, links, sub-tasks (up to 2000 chars).",
        "",
        "Press Add (or hit Enter in the title field).",
    ], Inches(0.5), Inches(1.3), Inches(5.3), Inches(5.5))
    add_placeholder_picture(s, prs,2, "The add form with title, client, date, time, and + Details expanded",
                            Inches(6.0), Inches(1.3), Inches(6.83), Inches(5.5))
    add_footer(s, prs, "Day Reminders v1.4.6")

    # 4. Three ways to add
    s = prs.slides.add_slide(blank)
    add_title_bar(s, prs, "Three ways to add")
    add_body_bullets(s, prs, [
        "1.  From the tab — the add form covered on the previous slide.",
        "",
        "2.  From bot chat — type slash commands in your Day Reminders chat:",
        "        /add 5pm tomorrow #work Send weekly report",
        "        /list  —  see what's open today",
        "        /done <substring>  —  mark a matching item done",
        "        /help  —  show all commands",
        "",
        "3.  From any Teams chat — click the ... menu under the message box,",
        "        pick Day Reminders → Quick add reminder, type, submit.",
        "        Great when you're mid-conversation and remember something.",
    ], Inches(0.6), Inches(1.3), Inches(12.0), Inches(5.5))
    add_footer(s, prs, "Day Reminders v1.4.6")

    # 5. Lines view
    s = prs.slides.add_slide(blank)
    add_title_bar(s, prs, "Lines view — your daily list")
    add_body_bullets(s, prs, [
        "Default view. One row per reminder, sorted by time.",
        "",
        "On each row you see:",
        "•  Title (prefixed with [Client] if a client is set)",
        "•  Colored chips for tags",
        "•  Outlined dashed chip for the client (click to filter)",
        "•  Time, with a custom lead-time badge if set",
        "•  ⋯ menu for row options (lead time + details textarea)",
        "",
        "Click any title, date, time, or chip to inline-edit.",
        "Enter saves, Esc cancels.",
    ], Inches(0.5), Inches(1.3), Inches(5.3), Inches(5.5))
    add_placeholder_picture(s, prs,3, "Lines view with multiple reminders showing tags, clients, and a high-priority star",
                            Inches(6.0), Inches(1.3), Inches(6.83), Inches(5.5))
    add_footer(s, prs, "Day Reminders v1.4.6")

    # 6. Week view
    s = prs.slides.add_slide(blank)
    add_title_bar(s, prs, "Week view — am I bombarded or free?")
    add_body_bullets(s, prs, [
        "For the moment when you need to know what your whole week looks like.",
        "",
        "Mon–Sun grid by due date. Today's column highlighted.",
        "Timed items stack by time within each day; anytime items at the bottom.",
        "",
        "Click any empty day to start adding a reminder for that date.",
        "",
        "Prev / Today / Next to navigate weeks. Mini Day ⇄ Week switcher in the header so you don't have to go back to the top bar to flip.",
    ], Inches(0.5), Inches(1.3), Inches(5.3), Inches(5.5))
    add_placeholder_picture(s, prs,4, "Week view with today highlighted and reminders stacked under several days",
                            Inches(6.0), Inches(1.3), Inches(6.83), Inches(5.5))
    add_footer(s, prs, "Day Reminders v1.4.6")

    # 7. Tags / Clients / Group / Filter
    s = prs.slides.add_slide(blank)
    add_title_bar(s, prs, "Organize — tags, clients, group, filter")
    add_body_bullets(s, prs, [
        "Tags — type #work, #urgent, #qc directly in the title. Colored chips, click to filter.",
        "",
        "Clients — set per reminder. Title shows as [Client] prefix everywhere.",
        "    Colored per client (deterministic — same client = same color across rows).",
        "    Click chip to filter; Shift+click or right-click to inline-edit.",
        "",
        "Group toggle — cycles off → tag → client → off (shortcut: g).",
        "    Sections per tag or per client, each header colored to match.",
        "",
        "Quick filters — All / Timed / Anytime / Priority / Done pills above the list.",
        "    Search box top-right (shortcut: f) matches title, tags, and client.",
    ], Inches(0.5), Inches(1.3), Inches(5.3), Inches(5.5))
    add_placeholder_picture(s, prs,5, "Lines view with Group: client active — multiple client sections, colored headers",
                            Inches(6.0), Inches(1.3), Inches(6.83), Inches(5.5))
    add_footer(s, prs, "Day Reminders v1.4.6")

    # 8. Notifications
    s = prs.slides.add_slide(blank)
    add_title_bar(s, prs, "Notifications — the proactive card")
    add_body_bullets(s, prs, [
        "When a timed reminder is due (or N minutes before, configurable):",
        "•  A card lands in your Day Reminders chat.",
        "•  Buttons: Mark done · Snooze 15m · Snooze 1h · Tomorrow.",
        "•  Title shows as [Client] Title; description (if any) appears below.",
        "•  Teams Activity Feed (bell icon) also gets a notification.",
        "",
        "End-of-day check-in:",
        "•  At your configured EOD time the bot posts 'Are you done?' with any open items.",
        "•  Buttons to acknowledge or snooze 15 minutes.",
        "",
        "Items not done by their due date auto-roll forward to today with an overdue Nd badge",
        "(capped at 30 days so old backlog doesn't pile up).",
    ], Inches(0.5), Inches(1.3), Inches(5.3), Inches(5.5))
    add_placeholder_picture(s, prs,6, "A real proactive Adaptive Card with title, time, description, and the four buttons",
                            Inches(6.0), Inches(1.3), Inches(6.83), Inches(5.5))
    add_footer(s, prs, "Day Reminders v1.4.6")

    # 9. Tips + Roadmap
    s = prs.slides.add_slide(blank)
    add_title_bar(s, prs, "Tips, shortcuts, what's next")
    add_body_bullets(s, prs, [
        "Keyboard shortcuts:",
        "    /  focus the add field           f  focus search",
        "    g  cycle group mode              v  cycle views",
        "    ?  open the quick guide          Esc  clear filters / cancel edit",
        "",
        "Other handy bits:",
        "•  Select multiple — top-bar button, then bulk-done / delete / star a batch.",
        "•  Templates — + Templates pill for common reminders, one click to add.",
        "•  Undo delete — toast hangs around for 5 seconds after every delete.",
        "•  Settings (top-right ⚙) — EOD check-in time, default lead minutes, weekdays-only, theme.",
        "",
        "Coming next (v1.5): sharing. Assign reminders to teammates; set per-tag default share lists",
        "    so #QC auto-shares with the QC team without picking recipients each time.",
        "",
        "Bugs, asks, 'this is annoying' feedback → drop them in the team channel or ping Josh.",
    ], Inches(0.6), Inches(1.3), Inches(12.0), Inches(5.5))
    add_footer(s, prs, "Day Reminders v1.4.6")

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    print(f"Wrote {OUT} ({OUT.stat().st_size:,} bytes, {len(prs.slides)} slides)")


if __name__ == "__main__":
    build_deck()
